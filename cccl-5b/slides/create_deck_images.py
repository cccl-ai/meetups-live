#!/usr/bin/env python3
"""
Screenshot all HTML slides with Playwright CLI and assemble into a PPTX.

Usage:
    uv run --with python-pptx python create_deck_images.py
"""

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu

SLIDES_DIR = Path(__file__).parent

SLIDE_ORDER = [
    "vikram_pawar.html",
    "uk_ambassadors.html",
    "ambassador_programme.html",
    "agenda_overview.html",
    "jan_peer.html",
    "ruslans_zavackis.html",
    "daniel_buchele.html",
    "valera_iatsko.html",
    "fryderyk_benigni.html",
    "aris_manshor.html",
    "fawaz_shah.html",
    "talha_sheikh.html",
    "rhys_cazenove.html",
    "thank_you.html",
]


def screenshot_html(html_path: Path, png_path: Path):
    url = f"file://{html_path.resolve()}"
    subprocess.run(
        ["npx", "playwright", "screenshot", "--viewport-size=1280,720", url, str(png_path)],
        check=True, capture_output=True,
    )


def main():
    screenshots_dir = SLIDES_DIR / "_screenshots"
    screenshots_dir.mkdir(exist_ok=True)

    for html_name in SLIDE_ORDER:
        html_path = SLIDES_DIR / html_name
        png_path = screenshots_dir / html_name.replace(".html", ".png")
        if not html_path.exists():
            print(f"  SKIP: {html_name}")
            continue
        print(f"  Screenshot: {html_name}")
        screenshot_html(html_path, png_path)

    # Build PPTX — 16:9 widescreen
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for html_name in SLIDE_ORDER:
        png_path = screenshots_dir / html_name.replace(".html", ".png")
        if not png_path.exists():
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            str(png_path), Emu(0), Emu(0),
            prs.slide_width, prs.slide_height,
        )

    output = SLIDES_DIR / "cccl_5b_deck.pptx"
    prs.save(str(output))
    print(f"\nCreated: {output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
