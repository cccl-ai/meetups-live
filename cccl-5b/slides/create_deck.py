#!/usr/bin/env python3
"""
Create the full CCCL #5B slide deck as an editable .pptx.

Usage:
    uv run --with python-pptx --with "qrcode[pil]" python create_deck.py
"""

import io
from pathlib import Path

import qrcode
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# CCCL brand colours
BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xE8, 0xE4, 0xDD)
ACCENT = RGBColor(0xC4, 0x82, 0x5F)
GREY = RGBColor(0xA8, 0xA4, 0xA0)
BODY = RGBColor(0xD4, 0xD0, 0xC8)
BLUE = RGBColor(0x88, 0xBB, 0xFF)

EVENT_TAG = "CCCL #5B — Claude Code for Agentic Engineering"
EVENT_DATE = "8 April 2026"
SLIDES_DIR = Path(__file__).parent


def make_qr(url: str) -> io.BytesIO:
    qr = qrcode.QRCode(version=1, box_size=10, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="white", back_color="#1A1A2E")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_bar(slide):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = EVENT_TAG
    p.font.size = Pt(11)
    p.font.color.rgb = GREY
    p.font.bold = True

    tb2 = slide.shapes.add_textbox(Inches(10), Inches(0.3), Inches(2.5), Inches(0.4))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = EVENT_DATE
    p2.font.size = Pt(11)
    p2.font.color.rgb = GREY
    p2.alignment = PP_ALIGN.RIGHT


def add_bottom_bar(slide):
    tb = slide.shapes.add_textbox(Inches(11), Inches(6.9), Inches(2), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "cccl.ai"
    p.font.size = Pt(11)
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.RIGHT


def add_accent_line(slide, left, top, width):
    line = slide.shapes.add_shape(1, left, top, width, Emu(Pt(3).emu))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_speaker_slide(prs, name, subtitle, bio, quote, linkedin_url, photo_filename=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Name
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(7), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.65), Inches(7), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT

    add_accent_line(slide, Inches(0.6), Inches(2.15), Inches(3.5))

    # Bio
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(7), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = bio
    p.font.size = Pt(17)
    p.font.color.rgb = BODY
    p.line_spacing = Pt(26)

    # Quote
    if quote:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(7), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'"{quote}"'
        p.font.size = Pt(15)
        p.font.italic = True
        p.font.color.rgb = BLUE

    # Photo
    if photo_filename:
        photo_path = SLIDES_DIR / photo_filename
        if photo_path.exists():
            slide.shapes.add_picture(
                str(photo_path), Inches(8.5), Inches(0.9), Inches(4.2), Inches(3.1)
            )

    # QR code
    if linkedin_url:
        qr_img = make_qr(linkedin_url)
        slide.shapes.add_picture(qr_img, Inches(9.5), Inches(4.3), Inches(2.0), Inches(2.0))
        tb = slide.shapes.add_textbox(Inches(9.2), Inches(6.35), Inches(2.5), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "LinkedIn"
        p.font.size = Pt(10)
        p.font.color.rgb = GREY
        p.alignment = PP_ALIGN.CENTER


def add_vikram_slide(prs):
    add_speaker_slide(
        prs,
        name="Vikram Pawar",
        subtitle="Claude Community Ambassador — Official Anthropic\nFounder, CCCL",
        bio="Built CCCL from scratch — 3000+ registrations, 5 events, and counting. "
            "16+ years running RUMQ — now exclusively AI & Claude Code evangelism. "
            "Daily Claude Code power user.",
        quote='Nothing is so permanent as a temporary government program — Milton Friedman\n'
              '...he said government program, but he clearly meant my codebase.',
        linkedin_url="https://www.linkedin.com/in/vikrammpawar/",
        photo_filename="vikram.jpg",
    )


def add_ambassadors_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Claude Community Ambassadors — UK"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(12), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Official Anthropic"
    p.font.size = Pt(16)
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.CENTER

    # Three ambassadors
    ambassadors = [
        ("Vikram Pawar", "London", "vikram_ambassador.jpg"),
        ("Olivier Legris", "London", "olivier_legris.jpeg"),
        ("Max Tatton-Brown", "London", "max_tatton_brown.jpeg"),
    ]

    for i, (name, city, photo) in enumerate(ambassadors):
        x = Inches(1.5 + i * 3.8)
        photo_path = SLIDES_DIR / photo
        if photo_path.exists():
            slide.shapes.add_picture(str(photo_path), x, Inches(2.3), Inches(2.5), Inches(2.5))

        tb = slide.shapes.add_textbox(x - Inches(0.3), Inches(5.0), Inches(3.1), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = slide.shapes.add_textbox(x - Inches(0.3), Inches(5.45), Inches(3.1), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = city
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER


def add_programme_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Claude Community Ambassador Programme"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(12), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Appointed by Anthropic. ~60 ambassadors worldwide building local Claude communities."
    p.font.size = Pt(15)
    p.font.color.rgb = GREY

    items = [
        ("Claude Community Events",
         "Host regular meetups and workshops, funded and supported by Anthropic. Branded content, swag, and event materials provided."),
        ("Claude Community Hackathons",
         "Partner with nonprofits, local government, and civic organisations to build real solutions with Claude for public good."),
        ("Impact Labs",
         "Community-driven deep sessions focused on tangible outcomes and local impact."),
        ("Builders Council",
         "Invited product feedback programme. Real-world usage and community perspective directly shapes what Anthropic builds."),
        ("Code with Claude",
         "Anthropic's flagship conference. Ambassadors get guaranteed attendance and can present on the Builders Stage."),
    ]

    y = 2.2
    for title, desc in items:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(11), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.3), Inches(11), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = BODY

        y += 0.9


def add_agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Tonight's Agenda"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Left column — First half
    left_items = [
        ("WELCOME", None),
        ("18:15", "Vikram & Rob Hart"),
        ("FIRST HALF — 3 TALKS", None),
        ("18:25", "Jan Peer — Heap Today, Gone Tomorrow"),
        ("18:37", "Ruslans Zavackis — Supercharge Obsidian"),
        ("18:49", "Daniel Büchele — Figma + coder.com"),
        ("BREAK", None),
        ("19:01", "Networking & drinks"),
    ]

    y = 1.7
    for time, text in left_items:
        if text is None:
            # Section header
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(5), Inches(0.3))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = time
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.font.letter_spacing = Pt(2)
            y += 0.3
        else:
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(0.6), Inches(0.25))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = time
            p.font.size = Pt(12)
            p.font.color.rgb = GREY

            tb = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(4.5), Inches(0.25))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            y += 0.28

    # QR code for agenda
    qr_img = make_qr("https://cccl-ai.github.io/meetups-live/cccl-5b/agenda")
    slide.shapes.add_picture(qr_img, Inches(1.5), Inches(5.0), Inches(1.5), Inches(1.5))
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(6.5), Inches(2.0), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Full agenda — scan"
    p.font.size = Pt(10)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Right column — Second half
    right_items = [
        ("SECOND HALF — 3 TALKS", None),
        ("19:15", "Valera Iatsko — Full stack, multi-region"),
        ("19:27", "Aris Manshor ← Fryderyk Benigni"),
        ("", "  Context Engineering best practices"),
        ("19:39", "Talha Sheikh ← Fawaz Shah"),
        ("", "  Building Vector — enforcement layer"),
        ("CLOSING", None),
        ("19:51", "Networking & open mic lightning talks"),
        ("20:50", "Closing — Vikram & Rob Hart"),
        ("21:00", "Pub"),
    ]

    y = 1.7
    for time, text in right_items:
        if text is None:
            tb = slide.shapes.add_textbox(Inches(6.8), Inches(y), Inches(5), Inches(0.3))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = time
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.font.letter_spacing = Pt(2)
            y += 0.3
        else:
            if time:
                tb = slide.shapes.add_textbox(Inches(6.8), Inches(y), Inches(0.6), Inches(0.25))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = time
                p.font.size = Pt(12)
                p.font.color.rgb = GREY

            tb = slide.shapes.add_textbox(Inches(7.5), Inches(y), Inches(5), Inches(0.25))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = text.strip()
            p.font.size = Pt(12) if time else Pt(10)
            p.font.color.rgb = WHITE if time else GREY
            y += 0.28

    # Co-hosted box
    tb = slide.shapes.add_textbox(Inches(7.5), Inches(5.2), Inches(4), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CO-HOSTED WITH"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.letter_spacing = Pt(1)
    p2 = tf.add_paragraph()
    p2.text = "Rob Hart"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p3 = tf.add_paragraph()
    p3.text = "AI for Engineers London"
    p3.font.size = Pt(12)
    p3.font.color.rgb = GREY


def add_thank_you_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Stay connected with the community"
    p.font.size = Pt(18)
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.CENTER

    qr_items = [
        ("WhatsApp", "cccl.ai/whatsapp", "https://cccl.ai/whatsapp"),
        ("Slack", "cccl.ai/slack", "https://cccl.ai/slack"),
        ("Luma", "cccl.ai/luma", "https://cccl.ai/luma"),
        ("Website", "cccl.ai", "https://cccl.ai"),
    ]

    for i, (label, url_text, url) in enumerate(qr_items):
        x = Inches(1.0 + i * 3.0)
        qr_img = make_qr(url)
        slide.shapes.add_picture(qr_img, x, Inches(2.8), Inches(2.0), Inches(2.0))

        tb = slide.shapes.add_textbox(x - Inches(0.2), Inches(5.0), Inches(2.4), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER

        tb = slide.shapes.add_textbox(x - Inches(0.2), Inches(5.35), Inches(2.4), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = url_text
        p.font.size = Pt(11)
        p.font.color.rgb = GREY
        p.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Vikram
    add_vikram_slide(prs)

    # 2. UK Ambassadors
    add_ambassadors_slide(prs)

    # 3. Ambassador Programme
    add_programme_slide(prs)

    # 4. Agenda
    add_agenda_slide(prs)

    # 5-8. First half speakers (Rob's)
    add_speaker_slide(prs,
        name="Jan Peer",
        subtitle="Software Engineer",
        bio="Heap Today, Gone Tomorrow — debugging Node.js memory issues with Claude Code + Chrome DevTools MCP.",
        quote="",
        linkedin_url="",
        photo_filename=None,
    )

    add_speaker_slide(prs,
        name="Ruslans Zavackis",
        subtitle="CrowdStrike",
        bio="Supercharge your Obsidian with Claude Code — building a knowledge management system that actually works.",
        quote="",
        linkedin_url="https://www.linkedin.com/in/ruslanzavacky/",
        photo_filename=None,
    )

    add_speaker_slide(prs,
        name="Daniel Büchele",
        subtitle="Figma",
        bio="How Figma uses coder.com and autonomous agents — scaling AI-assisted development across teams.",
        quote="",
        linkedin_url="https://www.linkedin.com/in/danielbuechele/",
        photo_filename=None,
    )

    # 9. Valera Iatsko
    add_speaker_slide(prs,
        name="Valera Iatsko",
        subtitle="Google",
        bio="Claude Code for the entire stack — from monorepo to multi-region cloud deployment.",
        quote="",
        linkedin_url="",
        photo_filename=None,
    )

    # 10. Fryderyk Benigni (introducer)
    add_speaker_slide(prs,
        name="Fryderyk Benigni",
        subtitle="Introducer — AppsVortex",
        bio="Travelling from Cambridge. Introducing Aris Manshor.",
        quote="",
        linkedin_url="https://www.linkedin.com/in/fryderykbenigni/",
        photo_filename=None,
    )

    # 11. Aris Manshor
    add_speaker_slide(prs,
        name="Aris Manshor",
        subtitle="Founder @ Kuro Data AI",
        bio="20+ years in software development. Building Kuro Data AI — multiple AI agents running "
            "simultaneously for one user request. Context engineering in production.",
        quote="Stay hungry, have a passion to learn every day until you die, fail quickly, "
              "don't dwell, learn fast and enjoy the journey",
        linkedin_url="https://www.linkedin.com/in/aris-manshor-9544644/",
        photo_filename="aris_manshor.jpg",
    )

    # 12. Fawaz Shah (introducer)
    add_speaker_slide(prs,
        name="Fawaz Shah",
        subtitle="Introducer — Senior SWE, Bloomberg",
        bio="NHS hospital comparison app built with Claude Code. Introducing Talha Sheikh.",
        quote="",
        linkedin_url="https://www.linkedin.com/in/fawaz-shah/",
        photo_filename=None,
    )

    # 13. Talha Sheikh
    add_speaker_slide(prs,
        name="Talha Sheikh",
        subtitle="AI Engineer @ Checkout",
        bio="Building Vector — a deterministic coding harness that guides coding agents in real time. "
            "Daily Claude Code user building an agentic AI OS.",
        quote="Building the enforcement layer I wish Claude had",
        linkedin_url="https://www.linkedin.com/in/talha-sheikh-007",
        photo_filename=None,
    )

    # 14. Thank You
    add_thank_you_slide(prs)

    output = SLIDES_DIR / "cccl_5b_deck.pptx"
    prs.save(str(output))
    print(f"Created: {output}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
